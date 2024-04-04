import os
import fitz  # PyMuPDF, 用于处理PDF
from pptx import Presentation  # 用于处理PPTX
import spacy

# 加载英语模型
nlp = spacy.load('en_core_web_sm')
def plural_to_singular_spacy(words):
    singular_words = []
    for phrase in words:
        # 处理短语中的每个单词
        doc = nlp(phrase.lower())
        singular_phrase = ' '.join(token.lemma_ for token in doc)  # 对每个token进行词形还原，然后重新组合
        singular_words.append(singular_phrase)
    return singular_words
def search_pdf(file_path, keywords):
    results = set() # 使用集合避免重复的页码
    doc = fitz.open(file_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text()
        for keyword in keywords:
            if keyword.lower() in text.lower():  # 直接检查整个短语是否存在
                results.add(page_num + 1)
                # 移除break; 以允许在同一页寻找所有关键词或短语
    return list(results)  # 转换成列表以便于后续处理

def search_pptx(file_path, keywords):
    results = set() # 使用集合避免重复的页码
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for keyword in keywords:
                    if keyword.lower() in shape.text.lower():  # 直接检查整个短语是否存在
                        results.add(i + 1)
                        # 移除break; 以允许在同一页寻找所有关键词或短语
    return list(results)  # 转换成列表以便于后续处理



def search_slides(folder_path, keywords):
    results = {}
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        if filename.endswith(".pdf"):
            matched_pages = search_pdf(file_path, keywords)
            if matched_pages:
                results[filename] = matched_pages
        elif filename.endswith(".pptx"):
            matched_pages = search_pptx(file_path, keywords)
            if matched_pages:
                results[filename] = matched_pages
    return results


def main():
    folder_path = "/Users/alexchen/Desktop/ADV AI/project/slides/CS0441"
    keywords = ['Translation','Logic','Negation','Language','Practice']
    keywords_lower = [keyword.lower() for keyword in keywords]
    singular_keywords = plural_to_singular_spacy(keywords_lower)
    print(singular_keywords)

    results1 = search_slides(folder_path, singular_keywords)
    print(results1)



if __name__ == "__main__":
    main()
